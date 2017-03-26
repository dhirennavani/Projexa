import logging
from pptx import Presentation
import os
import subprocess
from random import randint

from flask import Flask, render_template

from flask_ask import Ask, statement, question, session

from os import listdir
from os.path import isfile, join


app = Flask(__name__)

ask = Ask(app, "/")

logging.getLogger("flask_ask").setLevel(logging.DEBUG)

presentation_url="/home/pi/Projexa_Files/"

pptno=0

filenames=[]

slideno=0   
@ask.launch


def new_game():

    welcome_msg = render_template('welcome')
    onlyfiles = [f for f in listdir(presentation_url) if isfile(join(presentation_url, f))]
    filenames=onlyfiles
    return question(welcome_msg)


@ask.intent("OpenIntent", default={'ppt_index':0}, convert={'ppt_index':int})
def open_presentation(ppt_index):
    """
    numbers = [randint(0, 9) for _ in range(3)]

    round_msg = render_template('round', numbers=numbers)

    session.attributes['numbers'] = numbers[::-1]  # reverse

    return question(round_msg)
    """
    global pptno
    pptno=ppt_index

    global slideno
    slideno=0
    global filenames
    print filenames
    print ppt_index
    presentation_url_file=presentation_url+filenames[int(ppt_index)-1]
    s="libreoffice --show "+presentation_url_file+" --norestore --nolockcheck"
    print s
    a=subprocess.Popen(s,shell=True)
    return statement("Opened ")   


@ask.intent("CurrentDetailsIntent")
def current_details_intent():
    global slideno
    global pptno
    return statement("Current Slide Number is "+str(slideno+1))

@ask.intent("ListIntent")

def list_presentations():       
    onlyfiles = [f for f in listdir(presentation_url) if isfile(join(presentation_url, f))]
    print onlyfiles
    msg=""
    index=1
    global filenames
    filenames=onlyfiles
    for afile in onlyfiles:
        msg=msg+"To open "+afile+" . Say open "+str(index)+". "
        index=index+1
    return question(msg)

@ask.intent("CloseIntent")
def close_presentation():
    global slideno
    slideno=0
    print "something"
    a=subprocess.Popen( "xdotool search --name 'Impress' key Escape", shell=True)
    return statement("closed")

@ask.intent("NextIntent", default={'no_slides_right':1}, convert={'no_slides_right':int})
def next_slide(no_slides_right): 
    print "right"
    print str(no_slides_right)
    i=0
    
    global slideno
    slideno=slideno+1
    s="xdotool search --name 'Impress' key --delay 1000 Right"
    for i in range(1,int(no_slides_right)):
        s=s+" Right"
        slideno=slideno+1
    a=subprocess.Popen( s, shell=True)
    return statement("Moved")

@ask.intent("PrevIntent", default={'no_slides_left':1}, convert={'no_slides_left':int})
def prev_slide(no_slides_left):
    print "left"
    print str(no_slides_left)
    i=0
    global slideno
    x=slideno-int(no_slides_left)
    if x>=0:
        s="xdotool search --name 'Impress' key --delay 1000 Left"
        slideno=slideno-1
        for i in range(1,int(no_slides_left)):
            s=s+" Left"
            slideno=slideno-1
        a=subprocess.Popen( s, shell=True)
        return statement("Moved")
    else:
        return statement("Cannot Move sorry")












@ask.intent("IntroIntent" , mapping={'action': 'Action'})

def next_round(action):
	print ("my action is ----> " + action)
	session.attributes['Action'] = action
	session.attributes['callFirst']=0
	session.attributes['wordP']=''
	session.attributes['type']='simple'
	session.attributes['difficult word']=[]
	session.attributes['attempt']=1
	session.attributes['slp']=1
	session.attributes['mlp']=1
	session.attributes['clp']=1
	session.attributes['rlp']=1
	session.attributes['dlp']=1
	return question('name of the file?')
	
@ask.intent("filenameIntent" , mapping={'filename': 'filename'})

def number_word(filename):
	session.attributes['filename'] = filename
	from pptx import Presentation
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Hello, ASU!"
	subtitle.text = "We are creating a interactive projector with presentation skills!"
	prs.save("/home/pi/Projexa_Files/"+filename+'.pptx')
	return question('Okay. I am gona '+ session.attributes['Action']+', a presentation with file name '+ session.attributes['filename']+ '. What Next ?' )

@ask.intent("SpeechLearnIntent" , mapping={'word': 'Word'})

def speech_round(word):
	return question("please say the content you want to add in the next slide")

@ask.intent("CreateSlideIntent" , mapping={'title': 'title', "bullet":"bullet"})

def speech_round(title,bullet):
        f = open("/home/pi/Projexa_Files/"+session.attributes['filename']+'.pptx', 'rb')
	prs = Presentation(f)
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = title
	tf = body_shape.text_frame
	tf.text = bullet
	prs.save("/home/pi/Projexa_Files/"+session.attributes['filename']+'.pptx')
	f.close
	return statement("done")



    
if __name__ == '__main__':

    app.run(debug=True)
